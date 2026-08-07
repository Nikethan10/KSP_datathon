"""
PRAHARI — KSP Datathon 2026
PowerPoint Presentation Generator
Team: fault line | Challenge 02 — AI-Driven Crime Analytics & Visualization Platform

Run:  python make_ppt.py
Output: PRAHARI_KSP_Datathon_2026.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt, Emu, Cm
import copy
from lxml import etree
import os

# ── Colour palette ─────────────────────────────────────────────────────────────
CHARCOAL    = RGBColor(0x15, 0x18, 0x1C)   # #15181c  slide background
CHARCOAL2   = RGBColor(0x1E, 0x22, 0x28)   # slightly lighter panel
CHARCOAL3   = RGBColor(0x26, 0x2C, 0x34)   # card background
GOLD        = RGBColor(0xC9, 0xA3, 0x5C)   # #c9a35c brass accent
GOLD_LIGHT  = RGBColor(0xD9, 0xBB, 0x7A)   # lighter gold for hover
OFF_WHITE   = RGBColor(0xF2, 0xF0, 0xEB)   # primary text
SLATE       = RGBColor(0x94, 0xA3, 0xB8)   # secondary text
SLATE_DIM   = RGBColor(0x64, 0x74, 0x8B)   # tertiary / muted text
GREEN_OK    = RGBColor(0x4A, 0xDE, 0x80)   # positive metric
RED_FLAG    = RGBColor(0xF8, 0x71, 0x71)   # flagged / warning
BORDER      = RGBColor(0x2E, 0x35, 0x3F)   # border lines

# ── Slide dimensions (16:9, 10 × 5.625 in) ────────────────────────────────────
W = Inches(13.333)   # widescreen 16:9
H = Inches(7.5)

# ── Asset paths ───────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH   = os.path.join(BASE, "frontend", "public", "brand", "lockup.png")
MARK_PATH   = os.path.join(BASE, "frontend", "public", "brand", "mark-256.png")

# Generated images (may be absent on first run without image generation step)
ARTIFACT_DIR = os.path.join(
    os.path.expandvars(r"%USERPROFILE%"),
    r".gemini\antigravity\brain\2ebe210f-d0f2-46d1-b1df-06fdc3df0f3a"
)

def _find_img(names):
    """Find an image file by trying multiple name patterns."""
    import glob
    for name in names:
        # Direct path
        if os.path.exists(name):
            return name
        # In artifact dir (with timestamp suffix)
        pattern = os.path.join(ARTIFACT_DIR, f"{name}*.png")
        matches = sorted(glob.glob(pattern))
        if matches:
            return matches[-1]
    return None

FLOW_IMG = _find_img(["process_flow_diagram"])
ARCH_IMG = _find_img(["architecture_diagram"])

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color=CHARCOAL):
    """Set a solid background fill on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h,
             fill=None, line_color=None, line_width_pt=0.75):
    """Add a plain rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    sf = shape.fill
    if fill is None:
        sf.background()
    else:
        sf.solid()
        sf.fore_color.rgb = fill
    ln = shape.line
    if line_color:
        ln.color.rgb = line_color
        ln.width = Pt(line_width_pt)
    else:
        ln.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=14, bold=False, color=OFF_WHITE,
             align=PP_ALIGN.LEFT, italic=False,
             font="IBM Plex Sans", wrap=True, spacing_after=0):
    """Add a text box with precise formatting."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing_after:
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        pPr.set(qn('a:spcAft'), str(spacing_after * 100))
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=12, bold=False, color=OFF_WHITE,
             align=PP_ALIGN.LEFT, italic=False, font="IBM Plex Sans",
             space_before=0, space_after=0):
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        pPr = p._p.get_or_add_pPr()
        spB = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spB, qn('a:spcPts'))
        spcPts.set('val', str(int(space_before * 100)))
    if space_after:
        pPr = p._p.get_or_add_pPr()
        spA = etree.SubElement(pPr, qn('a:spcAft'))
        spcPts = etree.SubElement(spA, qn('a:spcPts'))
        spcPts.set('val', str(int(space_after * 100)))
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def add_image_safe(slide, path, x, y, w, h):
    """Add image if it exists, otherwise add a placeholder rect."""
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        r = add_rect(slide, x, y, w, h, fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.5)
        tf = r.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "[Image]"
        run.font.color.rgb = SLATE_DIM
        run.font.size = Pt(11)

def add_divider(slide, y, x=Inches(0.8), w=None, color=BORDER):
    """Horizontal thin line."""
    if w is None:
        w = W - Inches(1.6)
    add_rect(slide, x, y, w, Pt(0.75), fill=color, line_color=None)

def slide_number_stamp(slide, num, total=None):
    """Small monospace slide number bottom-right."""
    label = f"{num:02d}" if total is None else f"{num:02d} / {total:02d}"
    add_text(slide, label,
             W - Inches(1.4), H - Inches(0.45),
             Inches(1.2), Inches(0.35),
             size=8, color=SLATE_DIM, align=PP_ALIGN.RIGHT,
             font="IBM Plex Mono")

def file_stamp(slide, label="KSP-DTN-2026 · CH-02"):
    add_text(slide, label,
             Inches(0.5), H - Inches(0.45),
             Inches(4), Inches(0.35),
             size=7.5, color=SLATE_DIM,
             font="IBM Plex Mono")

def section_stamp(slide, text, x=Inches(0.8), y=Inches(0.55)):
    add_text(slide, text, x, y, Inches(8), Inches(0.3),
             size=8, color=GOLD, font="IBM Plex Mono", bold=True)

def slide_title(slide, text, x=Inches(0.8), y=Inches(0.9),
                w=None, size=28, color=OFF_WHITE):
    if w is None:
        w = W - Inches(3.5)
    add_text(slide, text, x, y, w, Inches(0.8),
             size=size, bold=True, color=color)

def logo_top_right(slide, w=Inches(2.2)):
    """Place the PRAHARI lockup logo in the top-right corner."""
    if os.path.exists(LOGO_PATH):
        aspect = 450 / 340   # approximate lockup aspect ratio (w/h)
        h = w / aspect
        x = W - w - Inches(0.5)
        y = Inches(0.35)
        slide.shapes.add_picture(LOGO_PATH, x, y, w, h)
    else:
        add_text(slide, "PRAHARI",
                 W - Inches(2.8), Inches(0.35), Inches(2.5), Inches(0.5),
                 size=14, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

def mark_top_right(slide, size=Inches(0.55)):
    """Place the small shield mark in the top-right corner."""
    if os.path.exists(MARK_PATH):
        x = W - size - Inches(0.5)
        y = Inches(0.35)
        slide.shapes.add_picture(MARK_PATH, x, y, size, size)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def build_slide_01(prs):
    """Slide 1 — Team Details"""
    sl = blank_slide(prs)
    fill_bg(sl)

    # Gold left accent bar
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)

    # Top gold band
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)

    # Logo top-right
    logo_top_right(sl)

    # Section stamp
    section_stamp(sl, "01  ·  TEAM DETAILS", y=Inches(0.6))

    # Big challenge title
    add_text(sl, "AI-Driven Crime Analytics\n& Visualization Platform",
             Inches(0.8), Inches(1.1), Inches(7.8), Inches(1.6),
             size=30, bold=True, color=OFF_WHITE)

    # Gold underline
    add_rect(sl, Inches(0.8), Inches(2.65), Inches(5), Pt(2), fill=GOLD)

    # Problem statement box
    ps_box = add_rect(sl, Inches(0.8), Inches(2.85), Inches(8.5), Inches(0.9),
                      fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.5)
    tf = ps_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = "PROBLEM STATEMENT  ·  Challenge 02"
    r0.font.name = "IBM Plex Mono"
    r0.font.size = Pt(7.5)
    r0.font.color.rgb = GOLD
    r0.font.bold = True
    add_para(tf, "Current systems rely on siloed data and manual reporting, limiting advanced analytics and proactive policing capabilities.",
             size=11, color=SLATE, font="IBM Plex Sans", space_before=6)

    # Team info grid
    grid_top = Inches(3.95)
    col1_x = Inches(0.8)
    col2_x = Inches(5.0)
    row_h = Inches(0.55)

    team_data_left = [
        ("TEAM NAME",   "fault line"),
        ("TEAM LEADER", "Katir"),
        ("TEAM SIZE",   "5 members"),
    ]
    team_data_right = [
        ("CHALLENGE", "02 — AI-Driven Crime Analytics"),
        ("EVENT",     "KSP Datathon 2026"),
        ("DEPLOYED",  "Zoho Catalyst · Live"),
    ]

    for i, (label, value) in enumerate(team_data_left):
        y = grid_top + i * row_h
        add_text(sl, label, col1_x, y, Inches(2.2), Inches(0.28),
                 size=7, color=GOLD, font="IBM Plex Mono", bold=True)
        add_text(sl, value, col1_x, y + Inches(0.23), Inches(2.2), Inches(0.3),
                 size=13, bold=True, color=OFF_WHITE)

    for i, (label, value) in enumerate(team_data_right):
        y = grid_top + i * row_h
        add_text(sl, label, col2_x, y, Inches(4), Inches(0.28),
                 size=7, color=GOLD, font="IBM Plex Mono", bold=True)
        add_text(sl, value, col2_x, y + Inches(0.23), Inches(4), Inches(0.3),
                 size=13, bold=True, color=OFF_WHITE)

    # Members strip
    members_y = grid_top + 3 * row_h + Inches(0.15)
    add_rect(sl, Inches(0.8), members_y, Inches(8.5), Pt(1), fill=BORDER)
    members_y += Inches(0.22)
    add_text(sl, "TEAM MEMBERS", Inches(0.8), members_y, Inches(2), Inches(0.28),
             size=7, color=GOLD, font="IBM Plex Mono", bold=True)
    add_text(sl, "Nikethan  ·  Hari Nair  ·  Katir  ·  Dhikshitha  ·  Nihan",
             Inches(0.8), members_y + Inches(0.25), Inches(8.5), Inches(0.35),
             size=13, color=OFF_WHITE, bold=True)

    file_stamp(sl)
    slide_number_stamp(sl, 1)
    return sl


def build_slide_02(prs):
    """Slide 2 — Brief about the solution"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "02  ·  BRIEF ABOUT THE SOLUTION")
    add_text(sl, "PRAHARI  ·  ಪ್ರಹರಿ",
             Inches(0.8), Inches(0.88), Inches(9), Inches(0.65),
             size=28, bold=True, color=OFF_WHITE)
    add_text(sl, "\"the sentinel\"  —  AI-driven crime intelligence for Karnataka State Police",
             Inches(0.8), Inches(1.52), Inches(9), Inches(0.4),
             size=12, color=GOLD, italic=True)

    add_divider(sl, Inches(2.0))

    # Data scale box
    data_box = add_rect(sl, Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.75),
                        fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.5)
    tf = data_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = "1,674,732 FIR records  ·  2016 – 2024  ·  41 districts  ·  1,074 police stations"
    r.font.name = "IBM Plex Mono"
    r.font.size = Pt(11)
    r.font.color.rgb = GOLD
    r.font.bold = True

    # Four-layer cards
    layers = [
        ("SENSE",   "Where is crime actually clustering?",
         "Getis-Ord Gi* / LISA at p < 0.05 — statistical significance,\nnot a blurred heatmap"),
        ("PREDICT", "What is likely next, and who is behind it?",
         "LightGBM with near-repeat features · STL anomaly\ndetection · Louvain co-offending communities"),
        ("ACT",     "Where do I send the units I have?",
         "Maximal-coverage integer program (OR-Tools)\nwith a greedy fallback"),
        ("TRUST",   "Why should I believe it?",
         "SHAP attributions · calibration curve ·\nreporting-bias-adjusted fairness audit"),
    ]

    card_w = Inches(2.8)
    card_h = Inches(2.55)
    card_y = Inches(3.1)
    gap = Inches(0.22)
    start_x = Inches(0.8)

    for i, (layer, question, method) in enumerate(layers):
        cx = start_x + i * (card_w + gap)
        card = add_rect(sl, cx, card_y, card_w, card_h,
                        fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.75)
        # Layer label
        add_text(sl, layer, cx + Inches(0.18), card_y + Inches(0.18),
                 card_w - Inches(0.36), Inches(0.38),
                 size=15, bold=True, color=GOLD, font="IBM Plex Mono")
        # Question
        add_text(sl, question, cx + Inches(0.18), card_y + Inches(0.6),
                 card_w - Inches(0.36), Inches(0.6),
                 size=10.5, bold=True, color=OFF_WHITE, wrap=True)
        # Method
        add_text(sl, method, cx + Inches(0.18), card_y + Inches(1.2),
                 card_w - Inches(0.36), Inches(1.2),
                 size=9, color=SLATE, wrap=True)

    # Footer note
    add_text(sl, "Every screen ends in a decision · Every recommendation explains itself · Deployed live on Zoho Catalyst",
             Inches(0.8), H - Inches(0.75), Inches(11.7), Inches(0.4),
             size=9, color=SLATE_DIM, align=PP_ALIGN.CENTER, italic=True)

    file_stamp(sl)
    slide_number_stamp(sl, 2)
    return sl


def build_slide_03(prs):
    """Slide 3 — Opportunities"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "03  ·  OPPORTUNITIES")
    slide_title(sl, "What Sets PRAHARI Apart", size=26)
    add_divider(sl, Inches(1.85))

    opp_blocks = [
        ("A", "HOW DIFFERENT FROM EXISTING IDEAS",
         [
             "Most crime dashboards show heatmaps. PRAHARI uses Getis-Ord Gi* at p < 0.05 — statistical significance, so a merely busy area is visibly distinguished from a genuinely hot one.",
             "Two baselines: state-wide (DGP resource allocation) and district-normalised (so an SP sees local hotspots, not just Bengaluru).",
             "We don't stop at prediction — we prescribe, via an integer program.",
         ]),
        ("B", "HOW IT SOLVES THE PROBLEM",
         [
             "Scattered tables → one unified data model across 8 CCTNS tables.",
             "Reactive policing → next-week risk forecast + STL anomaly alerts.",
             "Opaque outputs → SHAP attributions, calibration curves, fairness audit.",
         ]),
        ("C", "UNIQUE SELLING PROPOSITION",
         [
             '"Everyone shows where crime was. PRAHARI shows where it\'s about to be, who\'s behind it, and hands you a printable patrol plan — with the reasoning attached."',
         ]),
    ]

    y = Inches(2.05)
    for (letter, heading, bullets) in opp_blocks:
        # Letter badge
        badge = add_rect(sl, Inches(0.8), y, Inches(0.42), Inches(0.42),
                         fill=GOLD, line_color=None)
        tf = badge.text_frame
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = letter
        r.font.name = "IBM Plex Mono"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = CHARCOAL

        # Heading
        add_text(sl, heading, Inches(1.35), y, Inches(11), Inches(0.4),
                 size=8.5, bold=True, color=GOLD, font="IBM Plex Mono")

        # Bullets
        bullet_y = y + Inches(0.42)
        for b in bullets:
            add_text(sl, f"• {b}", Inches(1.35), bullet_y, Inches(11.3), Inches(0.55),
                     size=10.5, color=SLATE if letter != "C" else OFF_WHITE,
                     italic=(letter == "C"), wrap=True)
            bullet_y += Inches(0.55)

        y = bullet_y + Inches(0.25)

    file_stamp(sl)
    slide_number_stamp(sl, 3)
    return sl


def build_slide_04(prs):
    """Slide 4 — Features"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "04  ·  LIST OF FEATURES")
    slide_title(sl, "10 Core Requirements + Roadmap", size=26)
    add_divider(sl, Inches(1.85))

    features = [
        ("FR-1",  "Unified crime data model",              "8 CCTNS tables into one normalised schema"),
        ("FR-2",  "Geospatial drill-down",                 "District → station, click-to-fly, MapLibre GL"),
        ("FR-3",  "Gi*/LISA hotspots, 20 crime types",     "p < 0.05 significance; state + district baselines"),
        ("FR-4",  "Spatio-temporal risk forecast",          "LightGBM, near-repeat features, PAI 10.63 @ 5%"),
        ("FR-5",  "Anomaly & emerging-pattern alerts",      "STL residuals, changepoint detection, spike-first sort"),
        ("FR-6",  "Co-offending network + disruption",      "341,803 nodes · modularity 0.978 · Louvain communities"),
        ("FR-7",  "Socio-economic overlay",                 "District-level contextual factors on map"),
        ("FR-8",  "Patrol optimizer + briefing sheet",      "ILP maximal-coverage + greedy · printable PDF brief"),
        ("FR-9",  "SHAP explainability",                   "Per-prediction feature attributions, plain-English labels"),
        ("FR-10", "Fairness / bias audit",                  "Reporting-bias-adjusted Gini 0.183 · 11 districts flagged"),
    ]

    col_w = Inches(6.15)
    row_h = Inches(0.52)
    start_y = Inches(2.05)
    cols = [Inches(0.8), Inches(7.05)]

    for i, (fr, name, desc) in enumerate(features):
        col_idx = i % 2
        row_idx = i // 2
        cx = cols[col_idx]
        cy = start_y + row_idx * row_h

        # FR tag
        tag = add_rect(sl, cx, cy + Inches(0.06), Inches(0.72), Inches(0.36),
                       fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.5)
        tf = tag.text_frame
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = fr
        r.font.name = "IBM Plex Mono"
        r.font.size = Pt(7.5)
        r.font.bold = True
        r.font.color.rgb = GOLD

        # Feature name
        add_text(sl, name, cx + Inches(0.80), cy + Inches(0.06),
                 col_w - Inches(1.0), Inches(0.28),
                 size=11, bold=True, color=OFF_WHITE)
        # Description
        add_text(sl, desc, cx + Inches(0.80), cy + Inches(0.30),
                 col_w - Inches(1.0), Inches(0.22),
                 size=8.5, color=SLATE)

    # Extra features row
    extras_y = start_y + 5 * row_h + Inches(0.1)
    add_rect(sl, Inches(0.8), extras_y, Inches(12.5), Pt(1), fill=BORDER)
    extras_y += Inches(0.18)
    add_text(sl, "PLUS:  Bilingual English / ಕನ್ನಡ  ·  Offender dossiers  ·  Live intelligence feed",
             Inches(0.8), extras_y, Inches(9), Inches(0.35),
             size=10, color=GOLD, font="IBM Plex Sans")

    # Roadmap
    add_text(sl, "ROADMAP →  FR-11  Ask Prahari (natural language) · Complete Catalyst integration · Live CCTNS ingestion",
             Inches(0.8), extras_y + Inches(0.38), Inches(12), Inches(0.35),
             size=9.5, color=SLATE, italic=True)

    file_stamp(sl)
    slide_number_stamp(sl, 4)
    return sl


def build_slide_05(prs):
    """Slide 5 — Process Flow"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "05  ·  PROCESS FLOW")
    slide_title(sl, "Pipeline Architecture", size=26)
    add_divider(sl, Inches(1.85))

    if FLOW_IMG and os.path.exists(FLOW_IMG):
        # Show generated diagram
        sl.shapes.add_picture(FLOW_IMG, Inches(0.8), Inches(2.0),
                               Inches(12.0), Inches(5.0))
    else:
        # Draw flow manually using shapes
        stages = [
            ("01", "KSP FIR\nDataset", "1.67M records"),
            ("02", "Python\nPipeline", "9 steps, ~20 min"),
            ("03", "Cached\nArtefacts", "JSON / GeoJSON"),
            ("04", "Zoho\nCatalyst", "Hosting"),
            ("05", "React\nConsole", "SENSE/PREDICT\nACT/TRUST"),
        ]
        box_w = Inches(2.1)
        box_h = Inches(1.5)
        box_y = Inches(2.3)
        gap = Inches(0.25)
        total_w = len(stages) * box_w + (len(stages) - 1) * gap
        start_x = (W - total_w) / 2

        for i, (num, title, sub) in enumerate(stages):
            bx = start_x + i * (box_w + gap)
            b = add_rect(sl, bx, box_y, box_w, box_h,
                         fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.75)
            add_text(sl, num, bx + Inches(0.1), box_y + Inches(0.1),
                     Inches(0.5), Inches(0.3),
                     size=7.5, color=GOLD, font="IBM Plex Mono", bold=True)
            add_text(sl, title, bx + Inches(0.15), box_y + Inches(0.38),
                     box_w - Inches(0.3), Inches(0.65),
                     size=11, bold=True, color=OFF_WHITE)
            add_text(sl, sub, bx + Inches(0.15), box_y + Inches(1.05),
                     box_w - Inches(0.3), Inches(0.45),
                     size=8.5, color=SLATE)
            # Arrow between boxes
            if i < len(stages) - 1:
                ax = bx + box_w + Inches(0.04)
                ay = box_y + box_h / 2 - Pt(1)
                add_rect(sl, ax, ay, gap - Inches(0.08), Pt(2), fill=GOLD)

        # Outcome branches below last box
        outcomes = [
            "Identify hotspot",
            "Forecast next-week risk",
            "Deploy patrol + print briefing",
        ]
        last_box_x = start_x + (len(stages) - 1) * (box_w + gap)
        branch_center = last_box_x + box_w / 2
        branch_top_y = box_y + box_h + Inches(0.15)
        add_rect(sl, branch_center - Pt(1), branch_top_y,
                 Pt(2), Inches(0.25), fill=GOLD)

        out_w = Inches(2.8)
        total_out_w = len(outcomes) * out_w + (len(outcomes) - 1) * Inches(0.2)
        out_start_x = branch_center - total_out_w / 2
        out_y = branch_top_y + Inches(0.25)

        # Horizontal connector
        add_rect(sl, out_start_x + out_w / 2, out_y - Pt(1),
                 total_out_w - out_w, Pt(2), fill=GOLD)

        for j, outcome in enumerate(outcomes):
            ox = out_start_x + j * (out_w + Inches(0.2))
            add_rect(sl, ox + out_w / 2 - Pt(1), out_y, Pt(2), Inches(0.2), fill=GOLD)
            oc = add_rect(sl, ox, out_y + Inches(0.2), out_w, Inches(0.55),
                          fill=CHARCOAL2, line_color=GOLD, line_width_pt=0.5)
            tf = oc.text_frame
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.CENTER
            r = p0.add_run()
            r.text = outcome
            r.font.name = "IBM Plex Sans"
            r.font.size = Pt(9.5)
            r.font.color.rgb = OFF_WHITE

    file_stamp(sl)
    slide_number_stamp(sl, 5)
    return sl


def build_slide_06(prs):
    """Slide 6 — Wireframes (2x2 grid of tab screenshots described)"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "06  ·  PROTOTYPE WIREFRAMES")
    slide_title(sl, "Four Operational Tabs", size=26)
    add_divider(sl, Inches(1.85))

    tabs = [
        ("SENSE",   "Statistical hotspot map — Gi* / LISA at p < 0.05.\n722 hot cells / 16,650 analysed. State + district baselines.\nDistrict click → drill-down + trend chart."),
        ("PREDICT", "Risk forecast surface — LightGBM, week-ahead.\nAnomaly feed (STL spikes). Crime network 3D graph.\nLouvain communities + gang disruption ranking."),
        ("ACT",     "Patrol deployment map — ILP-optimised positions.\nPer-patrol briefing cards with crime type chips.\nPrint button → clean black-on-white PDF briefing."),
        ("TRUST",   "Reliability diagram + SHAP attribution chart.\nFairness audit — Gini 0.183 · 11 districts flagged.\nSHAP sample explanations with ground-truth validation."),
    ]

    card_w = Inches(5.85)
    card_h = Inches(2.35)
    gap = Inches(0.25)
    positions = [
        (Inches(0.8),             Inches(2.05)),
        (Inches(0.8) + card_w + gap, Inches(2.05)),
        (Inches(0.8),             Inches(2.05) + card_h + gap),
        (Inches(0.8) + card_w + gap, Inches(2.05) + card_h + gap),
    ]

    for (cx, cy), (tab, desc) in zip(positions, tabs):
        c = add_rect(sl, cx, cy, card_w, card_h,
                     fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.75)
        # Tab label
        label_box = add_rect(sl, cx + Inches(0.15), cy + Inches(0.12),
                              Inches(1.1), Inches(0.35),
                              fill=GOLD, line_color=None)
        tf = label_box.text_frame
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = tab
        r.font.name = "IBM Plex Mono"
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = CHARCOAL

        add_text(sl, desc, cx + Inches(0.15), cy + Inches(0.58),
                 card_w - Inches(0.3), card_h - Inches(0.75),
                 size=9.5, color=SLATE, wrap=True)

        # Screen prompt icon
        add_text(sl, "Live at → prahari-60076064719.development.catalystserverless.in/app/",
                 cx + Inches(0.15), cy + card_h - Inches(0.32),
                 card_w - Inches(0.3), Inches(0.28),
                 size=7.5, color=GOLD, font="IBM Plex Mono")

    file_stamp(sl)
    slide_number_stamp(sl, 6)
    return sl


def build_slide_07(prs):
    """Slide 7 — Architecture Diagram"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "07  ·  ARCHITECTURE DIAGRAM")
    slide_title(sl, "Three-Tier System Architecture", size=26)
    add_divider(sl, Inches(1.85))

    if ARCH_IMG and os.path.exists(ARCH_IMG):
        sl.shapes.add_picture(ARCH_IMG, Inches(0.8), Inches(2.05),
                               Inches(12.0), Inches(5.0))
    else:
        # Draw architecture manually
        # Tier labels on left
        tier_labels = [
            (Inches(2.05), "PRESENTATION LAYER", "React 19 + TypeScript · MapLibre GL · deck.gl · Recharts · Tailwind v4", CHARCOAL3),
            (Inches(3.55), "ANALYTICS LAYER  (Python, offline)", "", CHARCOAL2),
            (Inches(5.6),  "DATA LAYER", "KSP FIR Dataset  ·  8 CCTNS Tables", CHARCOAL3),
        ]

        tier_h = Inches(1.1)
        tier_w = Inches(9.8)
        tier_x = Inches(0.8)

        for ty, tier_name, tier_sub, tier_fill in tier_labels:
            add_rect(sl, tier_x, ty, tier_w, tier_h,
                     fill=tier_fill, line_color=BORDER, line_width_pt=0.5)
            add_text(sl, tier_name, tier_x + Inches(0.2), ty + Inches(0.12),
                     tier_w - Inches(0.4), Inches(0.35),
                     size=8, bold=True, color=GOLD, font="IBM Plex Mono")
            if tier_sub:
                add_text(sl, tier_sub, tier_x + Inches(0.2), ty + Inches(0.48),
                         tier_w - Inches(0.4), Inches(0.5),
                         size=10, color=OFF_WHITE)

        # Analytics layer boxes (inside middle tier)
        analytics_boxes = [
            ("SENSE",   "libpysal\nesda Gi*"),
            ("PREDICT", "LightGBM\nSTL · networkx"),
            ("ACT",     "OR-Tools\nILP"),
            ("TRUST",   "SHAP\nFairness"),
        ]
        ab_w = Inches(2.25)
        ab_h = Inches(0.85)
        ab_y = Inches(3.65)
        ab_gap = Inches(0.12)
        ab_start = tier_x + Inches(0.2)
        for j, (aname, asub) in enumerate(analytics_boxes):
            abx = ab_start + j * (ab_w + ab_gap)
            ab = add_rect(sl, abx, ab_y, ab_w, ab_h,
                          fill=CHARCOAL, line_color=GOLD, line_width_pt=0.75)
            add_text(sl, aname, abx + Inches(0.12), ab_y + Inches(0.08),
                     ab_w - Inches(0.24), Inches(0.3),
                     size=10, bold=True, color=GOLD, font="IBM Plex Mono")
            add_text(sl, asub, abx + Inches(0.12), ab_y + Inches(0.38),
                     ab_w - Inches(0.24), Inches(0.45),
                     size=8.5, color=SLATE)

        # Arrows between tiers
        arrow_x = tier_x + tier_w / 2
        add_rect(sl, arrow_x - Pt(1), Inches(4.6), Pt(2), Inches(0.25), fill=GOLD)
        add_rect(sl, arrow_x - Pt(1), Inches(5.2), Pt(2), Inches(0.25), fill=GOLD)

        # Zoho Catalyst band on right
        cat_x = tier_x + tier_w + Inches(0.25)
        cat_y = Inches(2.05)
        cat_h = Inches(4.65)
        cat_w = Inches(1.85)
        add_rect(sl, cat_x, cat_y, cat_w, cat_h,
                 fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.75)
        add_text(sl, "ZOHO\nCATALYST",
                 cat_x + Inches(0.1), cat_y + Inches(0.1),
                 cat_w - Inches(0.2), Inches(0.7),
                 size=9, bold=True, color=GOLD, font="IBM Plex Mono", align=PP_ALIGN.CENTER)

        services = [
            ("Hosting",    "Active"),
            ("Functions",  "Architected"),
            ("Data Store", "Architected"),
            ("Stratus",    "Architected"),
            ("Cron",       "Architected"),
        ]
        sy = cat_y + Inches(0.85)
        for sname, sstatus in services:
            add_text(sl, sname, cat_x + Inches(0.12), sy,
                     cat_w - Inches(0.24), Inches(0.28),
                     size=8.5, color=OFF_WHITE)
            sc = GREEN_OK if sstatus == "Active" else SLATE
            add_text(sl, sstatus, cat_x + Inches(0.12), sy + Inches(0.22),
                     cat_w - Inches(0.24), Inches(0.22),
                     size=7.5, color=sc, font="IBM Plex Mono", italic=True)
            sy += Inches(0.68)

    file_stamp(sl)
    slide_number_stamp(sl, 7)
    return sl


def build_slide_08(prs):
    """Slide 8 — Technologies"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "08  ·  TECHNOLOGIES")
    slide_title(sl, "Technology Stack", size=26)
    add_divider(sl, Inches(1.85))

    tech_rows = [
        ("Frontend",      "React 19  ·  TypeScript  ·  Vite 8  ·  MapLibre GL  ·  deck.gl  ·  Recharts  ·  Tailwind v4"),
        ("Analytics",     "pandas  ·  libpysal + esda  ·  scikit-learn  ·  LightGBM  ·  statsmodels  ·  networkx  ·  SHAP"),
        ("Optimisation",  "OR-Tools / PuLP  (integer linear programming)"),
        ("Platform",      "Zoho Catalyst  (Hosting · Functions · Data Store · Stratus · Cron)"),
    ]

    row_y = Inches(2.1)
    row_h = Inches(0.8)

    for layer, stack in tech_rows:
        # Layer badge
        badge = add_rect(sl, Inches(0.8), row_y + Inches(0.05),
                          Inches(1.8), Inches(0.55),
                          fill=GOLD, line_color=None)
        tf = badge.text_frame
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = layer.upper()
        r.font.name = "IBM Plex Mono"
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = CHARCOAL

        # Stack items
        add_text(sl, stack, Inches(2.75), row_y + Inches(0.1),
                 Inches(10), Inches(0.55),
                 size=11.5, color=OFF_WHITE, wrap=True)

        add_divider(sl, row_y + row_h, x=Inches(0.8), w=Inches(12))
        row_y += row_h

    # Performance note
    note_y = row_y + Inches(0.25)
    note = add_rect(sl, Inches(0.8), note_y, Inches(12), Inches(0.75),
                    fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.5)
    tf = note.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r = p0.add_run()
    r.text = "PERFORMANCE NOTE  "
    r.font.name = "IBM Plex Mono"
    r.font.size = Pt(8)
    r.font.bold = True
    r.font.color.rgb = GOLD
    add_para(tf, "Console is code-split — initial JS payload is 288 kB, not 2.5 MB. Analytics are pre-computed offline and served as static artefacts, so there is no per-query compute cost at runtime.",
             size=10.5, color=SLATE, font="IBM Plex Sans", space_before=4)

    file_stamp(sl)
    slide_number_stamp(sl, 8)
    return sl


def build_slide_09(prs):
    """Slide 9 — Catalyst Services (SCORED)"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "09  ·  CATALYST SERVICES  ·  SCORED SECTION")
    add_text(sl, "Zoho Catalyst Integration", Inches(0.8), Inches(0.9),
             Inches(9), Inches(0.65), size=26, bold=True, color=OFF_WHITE)
    add_divider(sl, Inches(1.85))

    services = [
        ("Catalyst Hosting",    "Serves the React console, public marketing site and all precomputed analytics artefacts as static files",   "Deployed",   True),
        ("Catalyst Functions",  "Query API over cached analytics artefacts; officer feedback submission endpoint",                              "Architected", False),
        ("Catalyst Data Store", "Normalised case, offender, and officer-feedback tables; structured schema defined in BRD",                    "Architected", False),
        ("Catalyst Stratus",    "Object storage for hotspot GeoJSON surfaces (~84 MB) and trained model artefacts",                           "Architected", False),
        ("Catalyst Cron",       "Nightly recompute of hotspots, risk surface, anomaly feed and patrol assignments (footer shows last-run timestamp)", "Architected", False),
    ]

    col_headers = ["SERVICE", "HOW PRAHARI USES IT", "STATUS"]
    col_x = [Inches(0.8), Inches(3.0), Inches(11.6)]
    col_w = [Inches(2.1), Inches(8.5), Inches(1.6)]

    header_y = Inches(2.0)
    # Header row background
    add_rect(sl, Inches(0.8), header_y, Inches(12.5), Inches(0.4),
             fill=CHARCOAL3, line_color=BORDER, line_width_pt=0.5)

    for ci, (hdr, cx, cw) in enumerate(zip(col_headers, col_x, col_w)):
        add_text(sl, hdr, cx + Inches(0.1), header_y + Inches(0.05),
                 cw, Inches(0.32),
                 size=8, bold=True, color=GOLD, font="IBM Plex Mono")

    row_h = Inches(0.85)
    for ri, (svc, usage, status, is_live) in enumerate(services):
        ry = header_y + Inches(0.42) + ri * row_h
        bg = CHARCOAL2 if ri % 2 == 0 else CHARCOAL3
        add_rect(sl, Inches(0.8), ry, Inches(12.5), row_h,
                 fill=bg, line_color=BORDER, line_width_pt=0.25)

        # Service name
        add_text(sl, svc, col_x[0] + Inches(0.1), ry + Inches(0.1),
                 col_w[0], Inches(0.4),
                 size=10, bold=True, color=GOLD)

        # Usage
        add_text(sl, usage, col_x[1] + Inches(0.1), ry + Inches(0.08),
                 col_w[1], row_h - Inches(0.16),
                 size=9.5, color=OFF_WHITE, wrap=True)

        # Status badge
        sc = GREEN_OK if is_live else SLATE
        status_label = "DEPLOYED" if is_live else "ARCHITECTED"
        sbadge = add_rect(sl, col_x[2], ry + Inches(0.22),
                           Inches(1.5), Inches(0.38),
                           fill=CHARCOAL if not is_live else RGBColor(0x14, 0x38, 0x22),
                           line_color=sc, line_width_pt=0.75)
        tf = sbadge.text_frame
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = status_label
        r.font.name = "IBM Plex Mono"
        r.font.size = Pt(7.5)
        r.font.bold = True
        r.font.color.rgb = sc

    # Honesty note
    note_y = header_y + Inches(0.42) + len(services) * row_h + Inches(0.12)
    add_text(sl, "Only Catalyst Hosting is live and verifiable by a judge in ten seconds. All other services are fully specified (schema, API surface, Cron schedule) and ready to activate once Catalyst CLI credentials are connected.",
             Inches(0.8), note_y, Inches(12.5), Inches(0.55),
             size=9, color=SLATE_DIM, italic=True, wrap=True)

    file_stamp(sl)
    slide_number_stamp(sl, 9)
    return sl


def build_slide_10(prs):
    """Slide 10 — Estimated Cost"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "10  ·  ESTIMATED COST")
    slide_title(sl, "Cost Model — Prototype → Production", size=26)
    add_divider(sl, Inches(1.85))

    cost_items = [
        ("Prototype (current)",
         "Catalyst free tier covers hosting, static artefact delivery, and all active features. Analytics run offline on a local machine. No per-query compute cost.",
         "₹ 0 / month", GREEN_OK),
        ("Nightly Cron execution",
         "One Catalyst Cron job triggers the analytics pipeline recompute (~22 min) nightly. Estimated at Catalyst standard Cron pricing.",
         "₹ ~500 / month", GOLD),
        ("Data Store storage",
         "Normalised CCTNS tables + feedback rows. ~2 GB estimated for full Karnataka dataset. Standard Catalyst Data Store rates.",
         "₹ ~800 / month", GOLD),
        ("Stratus object storage",
         "GeoJSON surfaces + model artefacts. ~84 MB compressed. Minimal cost.",
         "₹ ~200 / month", GOLD),
        ("Functions (query API)",
         "Cached-artefact queries — very low invocation rate (only when an officer loads a tab). Well within free tier for a pilot.",
         "₹ ~0 / month", GREEN_OK),
    ]

    row_h = Inches(0.88)
    row_y = Inches(2.1)
    for item_name, desc, cost, cost_color in cost_items:
        card = add_rect(sl, Inches(0.8), row_y, Inches(12.5), row_h,
                        fill=CHARCOAL3, line_color=BORDER, line_width_pt=0.5)
        add_text(sl, item_name, Inches(1.0), row_y + Inches(0.08),
                 Inches(7.5), Inches(0.32),
                 size=11, bold=True, color=OFF_WHITE)
        add_text(sl, desc, Inches(1.0), row_y + Inches(0.38),
                 Inches(7.5), Inches(0.45),
                 size=9, color=SLATE, wrap=True)
        add_text(sl, cost, Inches(9.8), row_y + Inches(0.25),
                 Inches(3.2), Inches(0.38),
                 size=12, bold=True, color=cost_color,
                 align=PP_ALIGN.RIGHT, font="IBM Plex Mono")
        row_y += row_h + Inches(0.06)

    # Summary note
    add_rect(sl, Inches(0.8), row_y + Inches(0.1), Inches(12.5), Pt(1.5), fill=GOLD)
    add_text(sl, "At production scale the only recurring cost is nightly Cron execution and Data Store storage. Analytics are pre-computed, so there is no per-officer-request compute charge.",
             Inches(0.8), row_y + Inches(0.28), Inches(12.5), Inches(0.5),
             size=10, color=SLATE, italic=True, align=PP_ALIGN.CENTER, wrap=True)

    file_stamp(sl)
    slide_number_stamp(sl, 10)
    return sl


def build_slide_11(prs):
    """Slide 11 — Prototype Snapshots"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "11  ·  PROTOTYPE SNAPSHOTS")
    slide_title(sl, "Live Console — Four Operational Views", size=26)
    add_divider(sl, Inches(1.85))

    snapshots = [
        ("SENSE",   "722 hot cells / 16,650 analysed visible in readout. Click district to drill in."),
        ("PREDICT", "Risk surface + AI Prediction Insight panel + anomaly feed."),
        ("ACT",     "Patrol beats deployed + briefing sheet open with per-unit crime summary."),
        ("TRUST",   "Reliability diagram + SHAP global attribution + fairness panel."),
    ]

    card_w = Inches(5.85)
    card_h = Inches(2.35)
    gap = Inches(0.25)

    for idx, (tab, caption) in enumerate(snapshots):
        col = idx % 2
        row = idx // 2
        cx = Inches(0.8) + col * (card_w + gap)
        cy = Inches(2.05) + row * (card_h + gap)

        # Placeholder card styled as screenshot frame
        c = add_rect(sl, cx, cy, card_w, card_h,
                     fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.75)

        # Tab bar simulation
        add_rect(sl, cx, cy, card_w, Inches(0.3), fill=CHARCOAL2, line_color=None)

        # Tab label
        tb = add_rect(sl, cx + Inches(0.08), cy + Inches(0.04),
                       Inches(0.9), Inches(0.22),
                       fill=GOLD, line_color=None)
        tf = tb.text_frame
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = tab
        r.font.name = "IBM Plex Mono"
        r.font.size = Pt(7)
        r.font.bold = True
        r.font.color.rgb = CHARCOAL

        # URL bar
        add_text(sl, "prahari-60076064719.development.catalystserverless.in/app/",
                 cx + Inches(1.1), cy + Inches(0.04), card_w - Inches(1.2), Inches(0.22),
                 size=6.5, color=SLATE_DIM, font="IBM Plex Mono")

        # Screenshot placeholder area
        add_rect(sl, cx + Inches(0.1), cy + Inches(0.35),
                  card_w - Inches(0.2), card_h - Inches(0.75),
                  fill=CHARCOAL, line_color=BORDER, line_width_pt=0.25)

        # Icon or text in placeholder
        add_text(sl, f"[ {tab} SCREENSHOT ]",
                 cx + Inches(0.1), cy + Inches(0.9),
                 card_w - Inches(0.2), Inches(0.4),
                 size=10, color=SLATE_DIM, font="IBM Plex Mono",
                 align=PP_ALIGN.CENTER)
        add_text(sl, "Take screenshot at 1280×800 from live site",
                 cx + Inches(0.1), cy + Inches(1.2),
                 card_w - Inches(0.2), Inches(0.35),
                 size=8, color=SLATE_DIM, align=PP_ALIGN.CENTER, italic=True)

        # Caption
        add_text(sl, caption,
                 cx + Inches(0.12), cy + card_h - Inches(0.38),
                 card_w - Inches(0.24), Inches(0.33),
                 size=8.5, color=GOLD, italic=False)

    file_stamp(sl)
    slide_number_stamp(sl, 11)
    return sl


def build_slide_12(prs):
    """Slide 12 — Performance / Benchmarking"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "12  ·  PERFORMANCE & BENCHMARKING")
    slide_title(sl, "Measured Results — All Reproducible", size=26)
    add_divider(sl, Inches(1.85))

    metrics = [
        ("RRI @ 5%",           "1.27",     "vs a 41.94% status-quo baseline on pre-2024 data",               OFF_WHITE),
        ("PAI @ top 5%",      "10.63",    "53.1% of all crime falls inside 5% of the map area",              GREEN_OK),
        ("Patrol coverage",   "11.67%",   "ILP-optimised vs 9.87% volume-driven status quo",                 GREEN_OK),
        ("Coverage uplift",   "+18.2%",   "vs status quo — same unit count, zero extra cost",                GREEN_OK),
        ("ILP verification",  "11.72%",   "Greedy within 0.05 pts of optimal — fast heuristic is reliable",  SLATE),
        ("Network",           "341,803",  "offender nodes · 509,633 links · modularity 0.978",               OFF_WHITE),
        ("Fairness Gini",     "0.183",    "Low disparity — 11 districts flagged for reporting-bias review",   GOLD),
    ]

    col_x = [Inches(0.8), Inches(3.9), Inches(6.0)]
    col_w = [Inches(3.0), Inches(2.0), Inches(7.0)]
    col_hdr = ["METRIC", "VALUE", "MEANING"]

    hdr_y = Inches(2.05)
    add_rect(sl, Inches(0.8), hdr_y, Inches(12.5), Inches(0.38),
             fill=CHARCOAL3, line_color=BORDER, line_width_pt=0.5)
    for ci, (hdr, cx, cw) in enumerate(zip(col_hdr, col_x, col_w)):
        add_text(sl, hdr, cx + Inches(0.1), hdr_y + Inches(0.05),
                 cw, Inches(0.28),
                 size=8, bold=True, color=GOLD, font="IBM Plex Mono")

    row_h = Inches(0.65)
    for ri, (metric, value, meaning, val_color) in enumerate(metrics):
        ry = hdr_y + Inches(0.4) + ri * row_h
        bg = CHARCOAL2 if ri % 2 == 0 else CHARCOAL3
        add_rect(sl, Inches(0.8), ry, Inches(12.5), row_h,
                 fill=bg, line_color=BORDER, line_width_pt=0.25)
        add_text(sl, metric, col_x[0] + Inches(0.1), ry + Inches(0.16),
                 col_w[0], Inches(0.38), size=10, color=OFF_WHITE)
        add_text(sl, value, col_x[1] + Inches(0.1), ry + Inches(0.1),
                 col_w[1], Inches(0.45), size=16, bold=True, color=val_color,
                 font="IBM Plex Mono")
        add_text(sl, meaning, col_x[2] + Inches(0.1), ry + Inches(0.16),
                 col_w[2], Inches(0.38), size=9.5, color=SLATE, wrap=True)

    # Note about baseline
    note_y = hdr_y + Inches(0.4) + len(metrics) * row_h + Inches(0.18)
    add_text(sl, "Note: Coverage uplift is stated vs volume-driven status quo (patrols at highest-volume stations) — never vs random. The random baseline flatters and will not survive a follow-up question.",
             Inches(0.8), note_y, Inches(12.5), Inches(0.45),
             size=8.5, color=SLATE_DIM, italic=True, wrap=True)

    file_stamp(sl)
    slide_number_stamp(sl, 12)
    return sl


def build_slide_12a(prs):
    """Slide 12a — Operational Impact"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "12a  ·  OPERATIONAL IMPACT")
    slide_title(sl, "What +18.2% Actually Means", size=26)
    add_divider(sl, Inches(1.85))

    # Context box
    ctx = add_rect(sl, Inches(0.8), Inches(2.05), Inches(12.5), Inches(0.7),
                   fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.5)
    tf = ctx.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = "Same six patrol units.  Same shift.  Same fuel budget."
    r.font.name = "IBM Plex Sans"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = GOLD

    # Comparison cards
    comp_y = Inches(2.95)
    comp_h = Inches(1.8)
    comp_w = Inches(5.75)

    # Status quo card
    sq = add_rect(sl, Inches(0.8), comp_y, comp_w, comp_h,
                  fill=CHARCOAL3, line_color=BORDER, line_width_pt=0.75)
    add_text(sl, "STATUS QUO", Inches(1.0), comp_y + Inches(0.15),
             comp_w - Inches(0.4), Inches(0.35),
             size=9, bold=True, color=SLATE, font="IBM Plex Mono")
    add_text(sl, "9.87%", Inches(1.0), comp_y + Inches(0.5),
             comp_w - Inches(0.4), Inches(0.7),
             size=40, bold=True, color=SLATE, font="IBM Plex Mono")
    add_text(sl, "of predicted crime inside patrolled radius\n(units placed at highest-volume stations)",
             Inches(1.0), comp_y + Inches(1.22), comp_w - Inches(0.4), Inches(0.5),
             size=9.5, color=SLATE, wrap=True)

    # VS
    add_text(sl, "VS", Inches(6.65), comp_y + Inches(0.65),
             Inches(0.9), Inches(0.5),
             size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
             font="IBM Plex Mono")

    # PRAHARI card
    pr = add_rect(sl, Inches(7.55), comp_y, comp_w, comp_h,
                  fill=RGBColor(0x14, 0x38, 0x22), line_color=GREEN_OK, line_width_pt=0.75)
    add_text(sl, "PRAHARI", Inches(7.75), comp_y + Inches(0.15),
             comp_w - Inches(0.4), Inches(0.35),
             size=9, bold=True, color=GREEN_OK, font="IBM Plex Mono")
    add_text(sl, "11.67%", Inches(7.75), comp_y + Inches(0.5),
             comp_w - Inches(0.4), Inches(0.7),
             size=40, bold=True, color=GREEN_OK, font="IBM Plex Mono")
    add_text(sl, "of predicted crime inside patrolled radius\n(units placed by ILP risk optimisation)",
             Inches(7.75), comp_y + Inches(1.22), comp_w - Inches(0.4), Inches(0.5),
             size=9.5, color=SLATE, wrap=True)

    # Uplift callout
    up_y = comp_y + comp_h + Inches(0.25)
    up = add_rect(sl, Inches(0.8), up_y, Inches(12.5), Inches(0.65),
                  fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.5)
    tf = up.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = "+18.2% coverage uplift · ILP-verified at 11.72% (greedy within 0.05 pts) · Zero additional cost"
    r.font.name = "IBM Plex Mono"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = GREEN_OK

    # Closing statement
    add_text(sl,
             "Scaled across 37 districts with nightly recomputation, that is a materially different deployment posture for no new spending.",
             Inches(0.8), up_y + Inches(0.78), Inches(12.5), Inches(0.5),
             size=10.5, color=SLATE, italic=True, align=PP_ALIGN.CENTER, wrap=True)

    file_stamp(sl)
    slide_number_stamp(sl, 12)
    return sl


def build_slide_13(prs):
    """Slide 13 — Links"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "13  ·  LINKS")
    slide_title(sl, "Resources & Submission Links", size=26)
    add_divider(sl, Inches(1.85))

    links = [
        ("GitHub",     "[pending — repo not yet pushed]",
         "Repository with full source, Python pipeline, and React console"),
        ("Demo Video", "[YouTube / Drive link — unlisted]",
         "3-minute walkthrough of SENSE → PREDICT → ACT → TRUST"),
        ("Deployed",   "https://prahari-60076064719.development.catalystserverless.in/app/",
         "Live prototype on Zoho Catalyst · No login required"),
        ("Stack Page", "https://prahari-60076064719.development.catalystserverless.in/app/#/stack",
         "Catalyst services table — verifies integration claims"),
    ]

    item_y = Inches(2.1)
    for label, url, desc in links:
        item = add_rect(sl, Inches(0.8), item_y, Inches(12.5), Inches(1.0),
                        fill=CHARCOAL3, line_color=BORDER, line_width_pt=0.5)

        # Label badge
        badge = add_rect(sl, Inches(1.0), item_y + Inches(0.12),
                          Inches(1.5), Inches(0.32),
                          fill=GOLD, line_color=None)
        tf = badge.text_frame
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = label.upper()
        r.font.name = "IBM Plex Mono"
        r.font.size = Pt(8)
        r.font.bold = True
        r.font.color.rgb = CHARCOAL

        add_text(sl, url, Inches(2.65), item_y + Inches(0.08),
                 Inches(10), Inches(0.42),
                 size=11, bold=True, color=GOLD, font="IBM Plex Mono")
        add_text(sl, desc, Inches(2.65), item_y + Inches(0.52),
                 Inches(10), Inches(0.38),
                 size=9.5, color=SLATE)

        item_y += Inches(1.12)

    file_stamp(sl)
    slide_number_stamp(sl, 13)
    return sl


def build_slide_14(prs):
    """Slide 14 — Future Development"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "14  ·  FUTURE DEVELOPMENT")
    slide_title(sl, "Roadmap & Next Steps", size=26)
    add_divider(sl, Inches(1.85))

    roadmap = [
        ("FR-11", "Ask Prahari — Natural Language Querying",
         "Specified and scoped. Officers query the analytics in plain Kannada or English: \"Show me theft hotspots in Mysuru this week.\" LLM layer over cached artefacts — no raw data exposure."),
        ("INT-01", "Complete Catalyst Integration",
         "Activate Functions (query API), Data Store (case/offender/feedback tables), Stratus (artefact storage), and Cron (nightly recompute). Infrastructure is fully specified in BRD."),
        ("ING-01", "Live CCTNS Ingestion",
         "Replace the static FIR export with a live pull from the CCTNS API, triggering incremental hotspot and risk recomputes as new cases are registered."),
        ("RET-01", "Officer Feedback Loop",
         "Patrol briefing feedback (thumbs up/down, already logged to localStorage) feeds into model retraining. Closes the human-in-the-loop cycle specified in BRD."),
        ("OPT-01", "Multi-shift Scheduling",
         "Extend the ILP to cover multiple shifts and heterogeneous unit types (motorcycle, vehicle, foot patrol) with different coverage radii."),
    ]

    item_y = Inches(2.05)
    item_h = Inches(0.9)
    for fr, title, desc in roadmap:
        tag = add_rect(sl, Inches(0.8), item_y + Inches(0.12),
                        Inches(0.9), Inches(0.3),
                        fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.5)
        tf = tag.text_frame
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = fr
        r.font.name = "IBM Plex Mono"
        r.font.size = Pt(7.5)
        r.font.bold = True
        r.font.color.rgb = GOLD

        add_text(sl, title, Inches(1.82), item_y + Inches(0.08),
                 Inches(11), Inches(0.32),
                 size=11, bold=True, color=OFF_WHITE)
        add_text(sl, desc, Inches(1.82), item_y + Inches(0.4),
                 Inches(11), Inches(0.45),
                 size=9.5, color=SLATE, wrap=True)

        add_divider(sl, item_y + item_h, x=Inches(0.8), w=Inches(12.5))
        item_y += item_h + Inches(0.06)

    file_stamp(sl)
    slide_number_stamp(sl, 14)
    return sl


def build_slide_addA(prs):
    """Add A — Responsible AI"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "ADD-A  ·  RESPONSIBLE AI")
    slide_title(sl, "Responsible by Design", size=26)
    add_divider(sl, Inches(1.85))

    principles = [
        ("Explainability",
         "SHAP attributions on every prediction. A reliability (calibration) diagram shows whether a 70% risk score actually means 70%. No black-box outputs."),
        ("Reporting Bias",
         "More FIRs can mean more crime or more police reporting. The fairness audit corrects for known differences before flagging any district. Gini 0.183 · 11 districts flagged for review."),
        ("Privacy",
         "Identifiers are hashed at ingest. No case-level or personal data appears on any public-facing page. The console is behind an authenticated console shell."),
        ("Human Oversight",
         "Every patrol recommendation is a proposal an officer can accept, edit or reject. Feedback (thumbs up / down) is logged. The algorithm advises — the officer decides."),
    ]

    py = Inches(2.05)
    for pname, pdesc in principles:
        # Gold dot
        add_rect(sl, Inches(0.8), py + Inches(0.14), Inches(0.14), Inches(0.14),
                 fill=GOLD, line_color=None)
        add_text(sl, pname, Inches(1.1), py + Inches(0.05),
                 Inches(11.5), Inches(0.32),
                 size=11, bold=True, color=OFF_WHITE)
        add_text(sl, pdesc, Inches(1.1), py + Inches(0.38),
                 Inches(11.5), Inches(0.45),
                 size=10, color=SLATE, wrap=True)
        py += Inches(0.95)

    # THE LINE WE DO NOT CROSS — callout box
    callout_y = py + Inches(0.1)
    callout = add_rect(sl, Inches(0.8), callout_y, Inches(12.5), Inches(1.25),
                       fill=RGBColor(0x1A, 0x14, 0x07),
                       line_color=GOLD, line_width_pt=1.5)
    tf = callout.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = "THE LINE WE DO NOT CROSS"
    r.font.name = "IBM Plex Mono"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = GOLD
    add_para(tf,
             "PRAHARI predicts risk for areas and time windows, and analyses networks of people already on record. "
             "It does NOT predict crime for named individuals. "
             "There is no pre-crime score for a person anywhere in this system.",
             size=11.5, bold=True, color=OFF_WHITE, align=PP_ALIGN.CENTER,
             space_before=8)

    file_stamp(sl)
    slide_number_stamp(sl, 15)
    return sl


def build_slide_addB(prs):
    """Add B — Known Limitations"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "ADD-B  ·  KNOWN LIMITATIONS")
    slide_title(sl, "What We Do Not Claim", size=26)
    add_divider(sl, Inches(1.85))

    # Intro text
    add_text(sl,
             "Volunteering our limits before a judge finds them reads as rigour. Getting caught reads as overselling.",
             Inches(0.8), Inches(2.05), Inches(12.5), Inches(0.4),
             size=11, color=SLATE, italic=True)

    limits = [
        ("Geography",   "Case coordinates are not district-faithful — every district's median coordinate lands near Bengaluru. We report location at police-station jurisdiction, the most granular level the data honestly supports."),
        ("Identity",    "PersonID is a within-case ordinal, not a cross-case identity. All offender linkage uses OffenderID. Cross-case identity resolution is out of scope."),
        ("Timestamps",  "CrimeRegisteredDate has no time component. Temporal analysis uses IncidentFromDate. Sub-daily temporal patterns cannot be modelled from this dataset."),
        ("Outcomes",    "CaseStatusID is single-valued, so case outcomes (conviction, acquittal) cannot be modelled. The fairness audit uses chargesheet filing as a proxy."),
        ("Coverage",    "Karnataka only. No multi-state dimension. Generalisability across Indian states has not been tested."),
    ]

    lim_y = Inches(2.6)
    for lname, ldesc in limits:
        # Tag
        tag = add_rect(sl, Inches(0.8), lim_y + Inches(0.08),
                        Inches(1.4), Inches(0.28),
                        fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.5)
        tf = tag.text_frame
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = lname.upper()
        r.font.name = "IBM Plex Mono"
        r.font.size = Pt(7.5)
        r.font.bold = True
        r.font.color.rgb = GOLD

        add_text(sl, ldesc, Inches(2.35), lim_y + Inches(0.04),
                 Inches(11), Inches(0.52),
                 size=10, color=SLATE, wrap=True)
        lim_y += Inches(0.72)

    # Closing note
    add_text(sl,
             "Each of these is a real quirk we hit and worked around. The workarounds are documented in STATUS.md and code comments.",
             Inches(0.8), lim_y + Inches(0.2), Inches(12.5), Inches(0.45),
             size=9.5, color=SLATE_DIM, italic=True, wrap=True)

    file_stamp(sl)
    slide_number_stamp(sl, 16)
    return sl


def build_slide_addC(prs):
    """Add C — Team Contribution"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "ADD-C  ·  TEAM CONTRIBUTION")
    slide_title(sl, "Who Built What", size=26)
    add_divider(sl, Inches(1.85))

    members = [
        ("Nikethan",    "Full-stack development lead. Frontend architecture (React 19 + TypeScript + Vite 8). SENSE and PREDICT tab UI. MapLibre GL integration and 3D deck.gl layers."),
        ("Hari Nair",   "Analytics pipeline. LightGBM risk model, feature engineering (near-repeat, temporal lags). STL anomaly detection. Model evaluation and benchmark harness."),
        ("Katir",       "Team lead & system design. BRD/PRD architecture documents. ACT (patrol optimizer) and TRUST (SHAP + fairness audit) layers. Zoho Catalyst deployment."),
        ("Dhikshitha",  "Co-offending network analysis. Louvain community detection, gang disruption simulation, networkx/3d-force-graph integration. Socio-economic overlay research."),
        ("Nihan",       "Data engineering. CCTNS schema normalisation, libpysal/esda Gi* spatial analysis, GeoJSON optimisation pipeline, district boundary data processing."),
    ]

    mem_y = Inches(2.05)
    for name, contrib in members:
        name_box = add_rect(sl, Inches(0.8), mem_y, Inches(2.0), Inches(0.7),
                             fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.75)
        tf = name_box.text_frame
        tf.word_wrap = False
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = name
        r.font.name = "IBM Plex Sans"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = GOLD

        add_text(sl, contrib, Inches(3.0), mem_y + Inches(0.08),
                 Inches(10.2), Inches(0.55),
                 size=10.5, color=OFF_WHITE, wrap=True)

        add_divider(sl, mem_y + Inches(0.75), x=Inches(0.8), w=Inches(12.5))
        mem_y += Inches(0.85)

    # Git note
    note_y = mem_y + Inches(0.1)
    add_rect(sl, Inches(0.8), note_y, Inches(12.5), Inches(0.6),
             fill=CHARCOAL3, line_color=BORDER, line_width_pt=0.5)
    add_text(sl,
             "Development was carried out largely on a shared machine, so git history reflects a single committer for much of the codebase. "
             "The breakdown above reflects actual ownership of each layer.",
             Inches(1.0), note_y + Inches(0.08), Inches(12.1), Inches(0.5),
             size=9, color=SLATE_DIM, italic=True, wrap=True)

    file_stamp(sl)
    slide_number_stamp(sl, 17)
    return sl


def build_slide_addD(prs):
    """Add D — Try it in 60 Seconds (Appendix)"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, W, Inches(0.12), fill=GOLD)
    add_rect(sl, 0, 0, Inches(0.18), H, fill=GOLD)
    mark_top_right(sl)

    section_stamp(sl, "APPENDIX  ·  JUDGE GUIDE")
    slide_title(sl, "Try It in 60 Seconds", size=26)
    add_divider(sl, Inches(1.85))

    url_box = add_rect(sl, Inches(0.8), Inches(2.05), Inches(12.5), Inches(0.6),
                       fill=CHARCOAL3, line_color=GOLD, line_width_pt=0.75)
    tf = url_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = "prahari-60076064719.development.catalystserverless.in/app/"
    r.font.name = "IBM Plex Mono"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = GOLD

    steps = [
        ("01", "Open the demo link → click  Enter Console"),
        ("02", "SENSE — click any district to drill in; note  722 hot cells / 16,650 analysed  in the readout"),
        ("03", "PREDICT — view the risk forecast surface, then switch to  Crime Network  in the top-left toggle"),
        ("04", "ACT — change the district and unit count, then open the briefing sheet and hit  Print"),
        ("05", "TRUST — inspect the reliability diagram and the fairness audit panel"),
    ]

    step_y = Inches(2.85)
    for num, step in steps:
        badge = add_rect(sl, Inches(0.8), step_y + Inches(0.06),
                          Inches(0.48), Inches(0.48),
                          fill=GOLD, line_color=None)
        tf = badge.text_frame
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run()
        r.text = num
        r.font.name = "IBM Plex Mono"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = CHARCOAL

        add_text(sl, step, Inches(1.45), step_y + Inches(0.1),
                 Inches(11.5), Inches(0.42),
                 size=12, color=OFF_WHITE, wrap=True)

        step_y += Inches(0.68)

    # Toggle tip
    tip_y = step_y + Inches(0.2)
    tip = add_rect(sl, Inches(0.8), tip_y, Inches(12.5), Inches(0.7),
                   fill=RGBColor(0x1A, 0x14, 0x07), line_color=GOLD, line_width_pt=0.75)
    tf = tip.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = "Toggle  ಕ  in the header to switch to a full Kannada interface.  No login required."
    r.font.name = "IBM Plex Sans"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = GOLD

    file_stamp(sl)
    slide_number_stamp(sl, 18)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("PRAHARI — KSP Datathon 2026 — PPT Generator")
    print("=" * 60)

    prs = new_prs()

    builders = [
        ("Slide 01 — Team Details",          build_slide_01),
        ("Slide 02 — Solution Brief",         build_slide_02),
        ("Slide 03 — Opportunities",          build_slide_03),
        ("Slide 04 — Features",               build_slide_04),
        ("Slide 05 — Process Flow",           build_slide_05),
        ("Slide 06 — Wireframes",             build_slide_06),
        ("Slide 07 — Architecture",           build_slide_07),
        ("Slide 08 — Technologies",           build_slide_08),
        ("Slide 09 — Catalyst Services",      build_slide_09),
        ("Slide 10 — Estimated Cost",         build_slide_10),
        ("Slide 11 — Prototype Snapshots",    build_slide_11),
        ("Slide 12 — Benchmarking",           build_slide_12),
        ("Slide 12a — Operational Impact",    build_slide_12a),
        ("Slide 13 — Links",                  build_slide_13),
        ("Slide 14 — Future Development",     build_slide_14),
        ("Slide A — Responsible AI",          build_slide_addA),
        ("Slide B — Limitations",             build_slide_addB),
        ("Slide C — Team Contribution",       build_slide_addC),
        ("Slide D — Try in 60s (Appendix)",   build_slide_addD),
    ]

    for label, builder in builders:
        print(f"  Building: {label}")
        try:
            builder(prs)
        except Exception as e:
            print(f"    ERROR in {label}: {e}")
            import traceback
            traceback.print_exc()

    out_path = os.path.join(BASE, "PRAHARI_KSP_Datathon_2026.pptx")
    prs.save(out_path)
    print()
    print(f"  Saved -> {out_path}")
    print(f"  Slides: {len(prs.slides)}")
    print("=" * 60)
    print("Done.")


if __name__ == "__main__":
    main()
